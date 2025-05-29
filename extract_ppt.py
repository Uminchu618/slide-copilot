from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json

# 保存先ディレクトリ
output_dir = "pptx_output"
if not os.path.exists(output_dir):
    os.makedirs(output_dir)


def extract_text_and_images(pptx_path):
    prs = Presentation(pptx_path)
    slide_info = []

    for idx, slide in enumerate(prs.slides):
        slide_texts = []
        image_files = []

        # テキスト抽出
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:  # 空文字でなければ追加
                    slide_texts.append(text)

            # 画像抽出
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                img_bytes = img.blob

                img_ext = img.ext
                img_filename = f"slide{idx+1}_img{len(image_files)+1}.{img_ext}"
                img_path = os.path.join(output_dir, img_filename)

                # 画像を書き出す
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                image_files.append(img_path)

        slide_info.append(
            {
                "slide_number": idx + 1,
                "texts": slide_texts,
                "images": image_files,
            }
        )
    return slide_info


# 使い方
pptx_file = "example.pptx"  # ←パワポファイルのパス
slide_contents = extract_text_and_images(pptx_file)

# JSONファイルとして保存
json_path = os.path.join(output_dir, "slide_texts.json")
with open(json_path, "w", encoding="utf-8") as f:
    # 文字化けしないようensure_ascii=False、インデントつき
    json.dump(slide_contents, f, ensure_ascii=False, indent=2)

print(f"スライド情報を {json_path} に保存しました。")
